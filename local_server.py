#!/usr/bin/env python3
"""
DJAI Local Artifact Server
Generates Excel (.xlsx), PowerPoint (.pptx), and JSON snapshots
from Quick Pitch data. Saves to ~/Desktop/DJAI/pitches/[TICKER]/

Usage:  python3 ~/Desktop/DJAI/local_server.py
        (or use start_local.sh)

Listens on http://localhost:3001
Endpoints:
  GET  /health   — liveness check, returns library availability
  POST /generate — generate all artifacts from pitch JSON payload
"""

import http.server
import json
import os
import sys
import traceback
from datetime import date
from pathlib import Path

# ── Library availability checks ────────────────────────────────────────────────
try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False
    print('[DJAI] Warning: openpyxl not found — Excel generation disabled.')
    print('[DJAI] Install: pip3 install openpyxl')

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print('[DJAI] Warning: python-pptx not found — PowerPoint generation disabled.')
    print('[DJAI] Install: pip3 install python-pptx')

PITCHES_DIR = Path.home() / 'Desktop' / 'DJAI' / 'pitches'
PORT = 3001

# ── Color palette (matches DJAI Bloomberg dark theme) ──────────────────────────
NAVY        = '1F3864'
NAVY_LIGHT  = '2E4A7A'
GOLD        = 'C9A84C'
GOLD_LIGHT  = 'E8C97A'
WHITE       = 'FFFFFF'
OFF_WHITE   = 'E8E6E0'
DARK_BG     = '0A0A0B'
SURFACE     = '111113'
GREEN       = '4CAF7D'
RED         = 'E05C5C'
AMBER       = 'E09B3C'
BLUE_INPUT  = '2E74B5'   # blue = hardcoded inputs (industry convention)
YELLOW_FILL = 'FFF2CC'   # yellow = key assumptions

# ── Template selector — mirrors JS pitchTemplateSelector ──────────────────────
import re as _re

_BANK_TICKERS = {
    'GS','MS','JPM','BAC','C','WFC','BLK','BX','APO','KKR','ARES','CG','TPG',
    'STT','BK','SCHW','COIN','IBKR','RJF','JEF','LAZ','EVR','HLI','PJT','MUFG',
    'DB','UBS','HSBC','BCS','TD','BMO','RY','CM','FITB','RF','HBAN','CFG',
}
_BANK_SECTOR_RE = _re.compile(
    r'financial|bank|invest.*bank|brokerage|asset manag|capital market|broker.dealer',
    _re.IGNORECASE
)

def select_template(data: dict) -> str:
    """Return 'bank_financials' or 'generic_corporate' based on payload."""
    ticker = str(data.get('ticker') or '').upper().split()[0].replace('-','')
    sector = str(data.get('sector') or '')
    family = str(data.get('templateFamily') or '')
    if (
        ticker in _BANK_TICKERS
        or _BANK_SECTOR_RE.search(sector)
        or family == 'bank_financials'
    ):
        return 'bank_financials'
    return 'generic_corporate'

# ── HTTP handler ───────────────────────────────────────────────────────────────
class ArtifactHandler(http.server.BaseHTTPRequestHandler):

    def log_message(self, fmt, *args):
        pass  # suppress default access log; print selectively below

    def _send_json(self, status, payload):
        body = json.dumps(payload, default=str).encode('utf-8')
        self.send_response(status)
        self.send_header('Content-Type', 'application/json')
        self.send_header('Content-Length', str(len(body)))
        self.send_header('Access-Control-Allow-Origin', '*')
        self.end_headers()
        self.wfile.write(body)

    def do_OPTIONS(self):
        self.send_response(200)
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'POST, GET, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')
        self.end_headers()

    def do_GET(self):
        if self.path == '/health':
            self._send_json(200, {
                'status': 'ok',
                'version': '1.1',
                'openpyxl': HAS_OPENPYXL,
                'pptx': HAS_PPTX,
                'pitches_dir': str(PITCHES_DIR)
            })
        elif self.path == '/portfolio':
            portfolio = load_portfolio_disk()
            self._send_json(200, {'ok': True, 'portfolio': portfolio})
        else:
            self._send_json(404, {'error': 'Not found'})

    def do_POST(self):
        try:
            length = int(self.headers.get('Content-Length', 0))
            raw = self.rfile.read(length)
            data = json.loads(raw.decode('utf-8'))
        except Exception as exc:
            self._send_json(400, {'error': f'Bad request: {exc}'})
            return

        if self.path == '/generate':
            try:
                result = generate_artifacts(data)
                ticker = data.get('ticker', '?')
                ok_count = sum(1 for a in result.get('artifacts', {}).values() if a.get('ok'))
                print(f'[DJAI] {ticker}: {ok_count}/{len(result.get("artifacts", {}))} artifacts saved')
                self._send_json(200, result)
            except Exception as exc:
                self._send_json(500, {'error': str(exc), 'trace': traceback.format_exc()})
        elif self.path == '/portfolio':
            try:
                save_portfolio_disk(data)
                self._send_json(200, {'ok': True})
            except Exception as exc:
                self._send_json(500, {'error': str(exc)})
        else:
            self._send_json(404, {'error': 'Not found'})


# ── Portfolio disk persistence ──────────────────────────────────────────────────
PORTFOLIO_FILE = Path.home() / 'Desktop' / 'DJAI' / 'portfolio.json'

def load_portfolio_disk() -> dict:
    try:
        if PORTFOLIO_FILE.exists():
            with open(PORTFOLIO_FILE, 'r', encoding='utf-8') as fh:
                return json.load(fh)
    except Exception:
        pass
    return {'positions': [], 'watchlist': [], 'alertThreshold': 0.20}

def save_portfolio_disk(portfolio: dict):
    PORTFOLIO_FILE.parent.mkdir(parents=True, exist_ok=True)
    with open(PORTFOLIO_FILE, 'w', encoding='utf-8') as fh:
        json.dump(portfolio, fh, indent=2, default=str)
    print(f'[DJAI] Portfolio saved: {len(portfolio.get("positions", []))} positions, '
          f'{len(portfolio.get("watchlist", []))} watchlist')


# ── Folder + file utilities ─────────────────────────────────────────────────────
def ensure_dir(path: Path):
    path.mkdir(parents=True, exist_ok=True)
    if not path.exists():
        raise RuntimeError(f'Failed to create directory: {path}')

def validate_file(path: Path, min_bytes: int = 50):
    """Returns (ok: bool, info: str|int)"""
    if not path.exists():
        return False, 'File not found after write'
    size = path.stat().st_size
    if size < min_bytes:
        return False, f'File too small ({size} bytes) — likely corrupt'
    return True, size


# ── Top-level orchestrator ─────────────────────────────────────────────────────
def generate_artifacts(data: dict) -> dict:
    ticker  = (data.get('ticker') or 'UNKNOWN').strip().upper()
    today   = date.today().isoformat()
    results = {'ticker': ticker, 'date': today, 'artifacts': {}}

    # Create output directory first
    output_dir = PITCHES_DIR / ticker
    try:
        ensure_dir(output_dir)
        results['folder'] = str(output_dir)
    except Exception as exc:
        results['error'] = f'Cannot create folder {output_dir}: {exc}'
        return results

    # 1. JSON Snapshot — cheapest; validates folder pipeline
    snap_path = output_dir / f'{ticker}_snapshot_{today}.json'
    try:
        snapshot = build_snapshot(data, ticker, today)
        with open(snap_path, 'w', encoding='utf-8') as fh:
            json.dump(snapshot, fh, indent=2, default=str)
        ok, info = validate_file(snap_path, min_bytes=50)
        results['artifacts']['snapshot'] = {'path': str(snap_path), 'ok': ok, 'info': info}
    except Exception as exc:
        results['artifacts']['snapshot'] = {'ok': False, 'error': str(exc)}

    # 2. Excel DCF model
    xlsx_path = output_dir / f'{ticker}_DCF_{today}.xlsx'
    if HAS_OPENPYXL:
        try:
            generate_excel(data, ticker, today, xlsx_path)
            ok, info = validate_file(xlsx_path, min_bytes=4096)
            results['artifacts']['excel'] = {'path': str(xlsx_path), 'ok': ok, 'info': info}
        except Exception as exc:
            results['artifacts']['excel'] = {'ok': False, 'error': str(exc)}
    else:
        results['artifacts']['excel'] = {'ok': False, 'error': 'openpyxl not installed'}

    # 3. PowerPoint pitch deck
    pptx_path = output_dir / f'{ticker}_Pitch_{today}.pptx'
    if HAS_PPTX:
        try:
            generate_pptx(data, ticker, today, pptx_path)
            ok, info = validate_file(pptx_path, min_bytes=4096)
            results['artifacts']['pptx'] = {'path': str(pptx_path), 'ok': ok, 'info': info}
        except Exception as exc:
            results['artifacts']['pptx'] = {'ok': False, 'error': str(exc)}
    else:
        results['artifacts']['pptx'] = {'ok': False, 'error': 'python-pptx not installed'}

    return results


# ── JSON Snapshot ──────────────────────────────────────────────────────────────
def build_snapshot(data: dict, ticker: str, today: str) -> dict:
    km  = data.get('keyMetrics') or {}
    dcf = data.get('dcf') or {}
    et  = data.get('entityType') or {}
    cal = data.get('calendar') or {}
    qs  = data.get('qualityScore') or {}
    return {
        'ticker':           ticker,
        'company':          data.get('company', ''),
        'date':             today,
        'entityType':       et.get('entityLabel', ''),
        'qualityScore':     qs.get('score'),
        'qualityLabel':     qs.get('label'),
        'currentPrice':     data.get('currentPrice'),
        'targetPrice':      data.get('targetPrice'),
        'rating':           data.get('rating', ''),
        'upside':           data.get('upside', ''),
        'thesis':           data.get('thesis', ''),
        'businessOverview': data.get('businessOverview', ''),
        'financialAnalysis':data.get('financialAnalysis', ''),
        'valSummary':       data.get('valSummary', ''),
        'keyMetrics':       km,
        'dcf': {
            'wacc':          dcf.get('wacc'),
            'terminalGrowth':dcf.get('terminalGrowth'),
            'impliedPrice':  dcf.get('impliedPrice'),
        },
        'bull':       data.get('bull', []),
        'bear':       data.get('bear', []),
        'catalysts':  data.get('catalysts', []),
        'risks':      data.get('risks', ''),
        'comps':      data.get('comps', []),
        'earningsDate':    cal.get('nextEarningsDate') or cal.get('nextEarningsDateEstimate'),
        'earningsDaysAway': cal.get('nextEarningsDaysAway'),
        'qcThesisScore':   data.get('qcThesisScore'),
        'qcFlag':          data.get('qcFlag'),
        'conclusion':      data.get('conclusion', ''),
        'timestamp':       today,
        'source':          'DJAI Quick Pitch',
    }


# ── Excel helper styles ────────────────────────────────────────────────────────
def _font(bold=False, size=10, color=None, italic=False):
    kw = {'bold': bold, 'size': size, 'italic': italic}
    if color:
        kw['color'] = color
    return Font(**kw)

def _fill(hex_color):
    return PatternFill('solid', fgColor=hex_color)

def _align(h='left', v='center', wrap=False):
    return Alignment(horizontal=h, vertical=v, wrap_text=wrap)

def _border_thin():
    s = Side(style='thin', color='CCCCCC')
    return Border(left=s, right=s, top=s, bottom=s)

def _set(ws, row, col, value, bold=False, size=10, color=None, fill=None,
         align='left', wrap=False, italic=False, num_fmt=None):
    cell = ws.cell(row=row, column=col, value=value)
    cell.font      = _font(bold=bold, size=size, color=color, italic=italic)
    cell.alignment = _align(h=align, wrap=wrap)
    if fill:
        cell.fill  = _fill(fill)
    if num_fmt:
        cell.number_format = num_fmt
    return cell


# ── Excel DCF Model ────────────────────────────────────────────────────────────
def generate_excel(data: dict, ticker: str, today: str, out_path: Path):
    km       = data.get('keyMetrics') or {}
    dcf      = data.get('dcf') or {}
    comp     = data.get('comps') or []
    hist_fin = data.get('historicalFinancials') or []
    wacc     = float(dcf.get('wacc') or 0.09)
    tgr      = float(dcf.get('terminalGrowth') or 0.025)
    imp      = float(dcf.get('impliedPrice') or 0)
    cur      = float(data.get('currentPrice') or 0)
    template = select_template(data)

    wb = openpyxl.Workbook()

    # ── Sheet 1: Summary ───────────────────────────────────────────────────────
    ws1 = wb.active
    ws1.title = 'Summary'
    ws1.column_dimensions['A'].width = 26
    ws1.column_dimensions['B'].width = 22
    ws1.column_dimensions['C'].width = 22
    ws1.column_dimensions['D'].width = 22
    ws1.column_dimensions['E'].width = 22

    # Header banner
    ws1.merge_cells('A1:E1')
    _set(ws1, 1, 1, f'DJAI — Quick Pitch: {ticker}  ·  {data.get("company", "")}',
         bold=True, size=14, color=GOLD_LIGHT, fill=NAVY, align='center')
    ws1.row_dimensions[1].height = 28

    ws1.merge_cells('A2:E2')
    _set(ws1, 2, 1, f'Generated: {today}  ·  Rating: {data.get("rating","")}  ·  Target: ${data.get("targetPrice","")}  ·  Current: ${data.get("currentPrice","")}  ·  Upside: {data.get("upside","")}',
         size=9, color=OFF_WHITE, fill=NAVY_LIGHT, align='center')
    ws1.row_dimensions[2].height = 18

    # Investment thesis
    ws1.row_dimensions[4].height = 14
    _set(ws1, 4, 1, 'INVESTMENT THESIS', bold=True, size=9, color=GOLD, fill=SURFACE)
    ws1.merge_cells('A5:E6')
    _set(ws1, 5, 1, data.get('thesis', ''), size=10, italic=True, align='left', wrap=True, color=OFF_WHITE)
    ws1.row_dimensions[5].height = 36

    # Key Metrics grid header
    ws1.row_dimensions[8].height = 16
    for col, hdr in enumerate(['METRIC', 'VALUE', 'METRIC', 'VALUE', ''], start=1):
        _set(ws1, 8, col, hdr, bold=True, size=8, color=WHITE, fill=NAVY, align='center')

    if template == 'bank_financials':
        metrics_left  = [('Revenue', km.get('revenue','')), ('ROTCE', km.get('rotce','')),
                         ('CET1 Ratio', km.get('cet1','')), ('Efficiency Ratio', km.get('efficiencyRatio',''))]
        metrics_right = [('Net Margin', km.get('netMargin','')), ('P/TBV', km.get('ptbv','')),
                         ('P/E', km.get('pe','')), ('Dividend Yield', km.get('dividendYield',''))]
    else:
        metrics_left  = [('Revenue', km.get('revenue','')), ('Rev Growth', km.get('revenueGrowth','')),
                         ('Gross Margin', km.get('grossMargin','')), ('Op Margin', km.get('operatingMargin',''))]
        metrics_right = [('Net Margin', km.get('netMargin','')), ('FCF', km.get('fcf','')),
                         ('P/E', km.get('pe','')), ('Fwd P/E', km.get('forwardPe',''))]
    for i, (lbl, val) in enumerate(metrics_left):
        row = 9 + i
        _set(ws1, row, 1, lbl, size=9, color=OFF_WHITE)
        _set(ws1, row, 2, val, size=9, color=WHITE, align='right')
    for i, (lbl, val) in enumerate(metrics_right):
        row = 9 + i
        _set(ws1, row, 3, lbl, size=9, color=OFF_WHITE)
        _set(ws1, row, 4, val, size=9, color=WHITE, align='right')

    # Bull / Bear
    _set(ws1, 14, 1, 'BULL CASE', bold=True, size=9, color=GREEN, fill=SURFACE)
    _set(ws1, 14, 3, 'BEAR CASE', bold=True, size=9, color=RED, fill=SURFACE)
    bull = data.get('bull') or []
    bear = data.get('bear') or []
    for i in range(3):
        row = 15 + i
        ws1.row_dimensions[row].height = 28
        b_text = f'▸  {bull[i]}' if i < len(bull) else ''
        r_text = f'▾  {bear[i]}' if i < len(bear) else ''
        cell_b = ws1.cell(row=row, column=1, value=b_text)
        cell_b.font = _font(size=9, color=OFF_WHITE)
        cell_b.alignment = _align(wrap=True)
        ws1.merge_cells(f'A{row}:B{row}')
        cell_r = ws1.cell(row=row, column=3, value=r_text)
        cell_r.font = _font(size=9, color=OFF_WHITE)
        cell_r.alignment = _align(wrap=True)
        ws1.merge_cells(f'C{row}:D{row}')

    # Analysis blurbs
    row = 19
    for label, text in [('Business', data.get('businessOverview','')),
                         ('Financials', data.get('financialAnalysis','')),
                         ('Valuation', data.get('valSummary',''))]:
        ws1.row_dimensions[row].height = 30
        _set(ws1, row, 1, label.upper(), bold=True, size=9, color=GOLD, fill=SURFACE)
        ws1.merge_cells(f'B{row}:E{row}')
        cell = ws1.cell(row=row, column=2, value=text)
        cell.font = _font(size=9, color=OFF_WHITE)
        cell.alignment = _align(wrap=True)
        row += 1

    # DCF disclaimer note
    ws1.merge_cells(f'A{row}:E{row}')
    _set(ws1, row, 1,
         '⚠  DCF figures are AI-estimated approximations — not a full bottom-up DCF model.  '
         'Use the DCF sheet as a directional starting point only.',
         size=8, color=AMBER, fill=SURFACE, italic=True, align='left', wrap=True)
    ws1.row_dimensions[row].height = 24

    ws1.sheet_view.showGridLines = False

    # ── Sheet 2: Valuation Model (template-dispatched) ────────────────────────
    if template == 'bank_financials':
        # ── Bank Valuation: P/TBV analysis + DDM framework ────────────────────
        ws2 = wb.create_sheet('Bank Valuation')
        ws2.column_dimensions['A'].width = 32
        for col in 'BCDEFG':
            ws2.column_dimensions[col].width = 16

        ws2.merge_cells('A1:G1')
        _set(ws2, 1, 1, f'BANK VALUATION — {ticker}', bold=True, size=13,
             color=GOLD_LIGHT, fill=NAVY, align='center')
        ws2.row_dimensions[1].height = 26

        ws2.merge_cells('A2:G2')
        _set(ws2, 2, 1,
             '⚠  Traditional FCF DCF is not applicable for banks. '
             'Use P/TBV multiple analysis and DDM as primary valuation frameworks.',
             size=8, color=AMBER, fill=SURFACE, italic=True, align='center')
        ws2.row_dimensions[2].height = 16

        # Key bank metrics block
        _set(ws2, 4, 1, 'KEY BANK METRICS', bold=True, size=9, color=GOLD, fill=SURFACE)
        bank_items = [
            ('ROTCE',           km.get('rotce','')),
            ('CET1 Capital Ratio', km.get('cet1','')),
            ('Efficiency Ratio', km.get('efficiencyRatio','')),
            ('Net Margin',      km.get('netMargin','')),
            ('Dividend Yield',  km.get('dividendYield','')),
            ('Book Value / Share', km.get('bookValuePerShare','')),
        ]
        for i, (lbl, val) in enumerate(bank_items):
            row = 5 + i
            _set(ws2, row, 1, lbl, size=9, color=OFF_WHITE)
            cell = ws2.cell(row=row, column=2, value=str(val) if val else '—')
            cell.font      = _font(size=9, color=GOLD_LIGHT, bold=True)
            cell.alignment = _align(h='right')

        # P/TBV Analysis
        row = 13
        _set(ws2, row, 1, 'P/TBV MULTIPLE ANALYSIS', bold=True, size=9, color=GOLD, fill=SURFACE)
        try:
            ptbv_current = float(str(km.get('ptbv',0) or 0).replace('x',''))
        except (ValueError, TypeError):
            ptbv_current = 0.0
        ptbv_scenarios = [
            ('Bear (0.8x TBV)',   0.8),
            ('Base (1.0x TBV)',   1.0),
            ('Bull (1.3x TBV)',   1.3),
            ('Premium (1.6x TBV)', 1.6),
        ]
        _set(ws2, row + 1, 1, 'Scenario', bold=True, size=8, color=WHITE, fill=NAVY, align='center')
        _set(ws2, row + 1, 2, 'P/TBV Multiple', bold=True, size=8, color=WHITE, fill=NAVY, align='center')
        _set(ws2, row + 1, 3, 'vs. Current P/TBV', bold=True, size=8, color=WHITE, fill=NAVY, align='center')
        for i, (scen, mult) in enumerate(ptbv_scenarios):
            r = row + 2 + i
            delta = f'{((mult / ptbv_current - 1) * 100):+.0f}%' if ptbv_current else '—'
            is_base = (mult == 1.0)
            fill = NAVY_LIGHT if is_base else SURFACE
            c_color = GOLD if is_base else OFF_WHITE
            _set(ws2, r, 1, scen, size=9, color=c_color, fill=fill)
            cell = ws2.cell(row=r, column=2, value=f'{mult:.1f}x')
            cell.font = _font(size=9, color=GOLD_LIGHT if is_base else WHITE, bold=is_base)
            cell.fill = _fill(fill)
            cell.alignment = _align(h='right')
            _set(ws2, r, 3, delta, size=9, color=GREEN if not str(delta).startswith('-') else RED,
                 fill=fill, align='right')

        # DDM Framework note
        row = row + 2 + len(ptbv_scenarios) + 2
        ws2.merge_cells(f'A{row}:G{row}')
        _set(ws2, row, 1, 'DIVIDEND DISCOUNT MODEL (DDM) FRAMEWORK', bold=True,
             size=9, color=GOLD, fill=SURFACE)
        ddm_note = (
            'For banks, DDM values the firm based on sustainable dividend capacity. '
            'P₀ = D₁ / (Ke − g), where D₁ = next dividend, Ke = cost of equity, g = sustainable growth rate. '
            f'ROTCE of {km.get("rotce","—")} × retention ratio ≈ sustainable ROE-based growth. '
            f'Current dividend yield: {km.get("dividendYield","—")}. '
            'Refine with consensus DPS estimates for a bottom-up DDM.'
        )
        row += 1
        ws2.merge_cells(f'A{row}:G{row + 1}')
        cell = ws2.cell(row=row, column=1, value=ddm_note)
        cell.font = _font(size=9, color=OFF_WHITE, italic=True)
        cell.alignment = _align(wrap=True)
        ws2.row_dimensions[row].height = 32

        ws2.sheet_view.showGridLines = False

    else:
        # ── Corporate DCF Model ────────────────────────────────────────────────
        ws2 = wb.create_sheet('DCF Model')
        ws2.column_dimensions['A'].width = 30
        for col in 'BCDEFG':
            ws2.column_dimensions[col].width = 16

        ws2.merge_cells('A1:G1')
        _set(ws2, 1, 1, f'DCF MODEL — {ticker}', bold=True, size=13,
             color=GOLD_LIGHT, fill=NAVY, align='center')
        ws2.row_dimensions[1].height = 26

        ws2.merge_cells('A2:G2')
        _set(ws2, 2, 1,
             '⚠  Simplified projection — not a full DCF model.  '
             'AI-estimated assumptions; verify before use.',
             size=8, color=AMBER, fill=SURFACE, italic=True, align='center')
        ws2.row_dimensions[2].height = 16

        # Assumptions block (yellow = user-editable inputs)
        _set(ws2, 3, 1, 'KEY ASSUMPTIONS', bold=True, size=9, color=GOLD, fill=SURFACE)
        assumptions = [
            ('WACC',                    wacc, '0.00%'),
            ('Terminal Growth Rate',    tgr,  '0.00%'),
            ('Implied Price (DCF)',      imp,  '"$"#,##0.00'),
            ('Current Price',           cur,  '"$"#,##0.00'),
            ('DCF Upside / (Downside)', (imp - cur) / cur if cur else 0, '0.0%'),
        ]
        for i, (lbl, val, fmt) in enumerate(assumptions):
            row = 4 + i
            _set(ws2, row, 1, lbl, size=9, color=OFF_WHITE)
            cell = ws2.cell(row=row, column=2, value=val)
            cell.font      = _font(size=9, color=BLUE_INPUT, bold=True)
            cell.fill      = _fill(YELLOW_FILL)
            cell.alignment = _align(h='right')
            cell.number_format = fmt

        # Revenue base: use historicalFinancials most recent year when available
        rev_base = 0.0
        rev_growth = 0.06
        if hist_fin:
            try:
                hf = hist_fin[0]
                rev_base   = float(hf.get('revenue') or 0)
                prev_rev   = float((hist_fin[1] if len(hist_fin) > 1 else {}).get('revenue') or 0)
                if prev_rev:
                    rev_growth = (rev_base - prev_rev) / prev_rev
            except (ValueError, TypeError, IndexError):
                pass
        if not rev_base:
            rev_base_str = str(km.get('revenue', '') or '').replace('$','').replace('B','').replace('M','').strip()
            try:
                rev_multiplier = 1e9 if 'B' in str(km.get('revenue','')) else 1e6
                rev_base = float(rev_base_str) * rev_multiplier
            except (ValueError, TypeError):
                rev_base = 0.0
        if not rev_growth:
            rev_growth_str = str(km.get('revenueGrowth', '') or '').replace('%','').strip()
            try:
                rev_growth = float(rev_growth_str) / 100
            except (ValueError, TypeError):
                rev_growth = 0.06
        op_margin_str = str(km.get('operatingMargin', '') or '').replace('%','').strip()
        try:
            op_margin = float(op_margin_str) / 100
        except (ValueError, TypeError):
            op_margin = 0.15
        # Clamp to sane range
        rev_growth = max(-0.30, min(rev_growth, 0.60))
        op_margin  = max(0.0,   min(op_margin,  0.60))

        row = 11
        _set(ws2, row, 1, 'REVENUE & FCF PROJECTIONS', bold=True, size=9, color=GOLD, fill=SURFACE)
        hdrs = ['Metric'] + [str(2025 + i) for i in range(5)]
        for c, h in enumerate(hdrs, start=1):
            _set(ws2, row + 1, c, h, bold=True, size=8, color=WHITE, fill=NAVY, align='center')
        ws2.row_dimensions[row + 1].height = 15

        years_data = []
        rev = rev_base
        for i in range(5):
            rev  = rev * (1 + rev_growth)
            fcf  = rev * op_margin * 0.7   # op income × (1 − tax) as FCF proxy
            df   = 1 / ((1 + wacc) ** (i + 1))
            pv   = fcf * df
            years_data.append({'rev': rev, 'fcf': fcf, 'df': df, 'pv': pv})

        # Show historical FCF anchor if available
        if hist_fin:
            hist_label_row = row - 1
            ws2.merge_cells(f'A{hist_label_row}:G{hist_label_row}')
            hist_note = 'Historical anchor (Alpha Vantage): ' + '  |  '.join(
                f'{h.get("year","?")}: Rev ${h.get("revenue",0)/1e9:.1f}B  FCF ${(h.get("freeCashflow") or 0)/1e9:.1f}B'
                for h in hist_fin[:3] if h.get('revenue')
            )
            _set(ws2, hist_label_row, 1, hist_note, size=7, color=GOLD, italic=True, fill=SURFACE)

        proj_rows = [
            ('Revenue ($B)',    [f'{y["rev"]/1e9:.2f}' for y in years_data]),
            ('FCF ($B)',        [f'{y["fcf"]/1e9:.2f}' for y in years_data]),
            ('Discount Factor', [f'{y["df"]:.4f}'      for y in years_data]),
            ('PV of FCF ($B)',  [f'{y["pv"]/1e9:.2f}'  for y in years_data]),
        ]
        for ri, (lbl, vals) in enumerate(proj_rows):
            r = row + 2 + ri
            _set(ws2, r, 1, lbl, size=9, color=OFF_WHITE)
            for ci, v in enumerate(vals, start=2):
                _set(ws2, r, ci, v, size=9, color=WHITE, align='right')

        total_pv = sum(y['pv'] for y in years_data)
        last_fcf = years_data[-1]['fcf'] * (1 + tgr)
        term_val = last_fcf / (wacc - tgr) if (wacc - tgr) > 0 else 0
        pv_term  = term_val / ((1 + wacc) ** 5)
        r = row + 6
        for lbl, val, fmt in [
            ('Sum of PV (FCF)',        total_pv,           '"$"#,##0,,,"B"'),
            ('Terminal Value',         term_val,           '"$"#,##0,,,"B"'),
            ('PV of Terminal Value',   pv_term,            '"$"#,##0,,,"B"'),
            ('Total Enterprise Value', total_pv + pv_term, '"$"#,##0,,,"B"'),
        ]:
            _set(ws2, r, 1, lbl, size=9, color=OFF_WHITE)
            cell = ws2.cell(row=r, column=2, value=round(val, 2))
            cell.font = _font(size=9, color=GOLD_LIGHT, bold=True)
            cell.alignment = _align(h='right')
            cell.number_format = fmt
            r += 1

        # Sensitivity table: WACC × TGR → implied price
        row = r + 2
        ws2.merge_cells(f'A{row}:G{row}')
        _set(ws2, row, 1, 'SENSITIVITY — Implied Price vs. WACC × Terminal Growth Rate',
             bold=True, size=9, color=GOLD, fill=SURFACE)
        row += 1
        wacc_steps = [-0.02, -0.01, 0, 0.01, 0.02]
        tgr_steps  = [-0.005, 0, 0.005, 0.010, 0.015]
        _set(ws2, row, 1, 'WACC \\ TGR', bold=True, size=8, color=WHITE, fill=NAVY, align='center')
        for ci, ts in enumerate(tgr_steps, start=2):
            _set(ws2, row, ci, f'{(tgr+ts)*100:.1f}%', bold=True, size=8, color=WHITE, fill=NAVY, align='center')
        for ri, ws in enumerate(wacc_steps):
            r2 = row + 1 + ri
            _set(ws2, r2, 1, f'{(wacc+ws)*100:.1f}%', bold=True, size=8, color=GOLD_LIGHT, fill=NAVY_LIGHT)
            for ci, ts in enumerate(tgr_steps, start=2):
                adj_wacc = wacc + ws
                adj_tgr  = tgr  + ts
                if adj_wacc > adj_tgr and adj_wacc > 0:
                    base_spread = wacc - tgr
                    new_spread  = adj_wacc - adj_tgr
                    adj_price   = imp * base_spread / new_spread if new_spread > 0 else 0
                else:
                    adj_price = 0
                is_base = (ws == 0 and ts == 0)
                cell = ws2.cell(row=r2, column=ci, value=round(adj_price, 2))
                cell.font   = _font(size=9, color=NAVY if is_base else OFF_WHITE, bold=is_base)
                cell.fill   = _fill(GOLD if is_base else SURFACE)
                cell.alignment = _align(h='center')
                cell.number_format = '"$"#,##0.00'

        ws2.sheet_view.showGridLines = False

    # ── Sheet 3: Comps ─────────────────────────────────────────────────────────
    ws3 = wb.create_sheet('Comps')
    ws3.column_dimensions['A'].width = 28
    ws3.column_dimensions['B'].width = 10
    ws3.column_dimensions['C'].width = 14
    ws3.column_dimensions['D'].width = 12
    ws3.column_dimensions['E'].width = 12
    ws3.column_dimensions['F'].width = 16

    ws3.merge_cells('A1:F1')
    _set(ws3, 1, 1, f'COMPARABLE COMPANIES — {ticker}', bold=True, size=13,
         color=GOLD_LIGHT, fill=NAVY, align='center')
    ws3.row_dimensions[1].height = 26

    if template == 'bank_financials':
        comp_hdrs = ['Company', 'Ticker', 'P/TBV', 'P/E', 'ROTCE', 'Premium / Discount']
    else:
        comp_hdrs = ['Company', 'Ticker', 'EV/EBITDA', 'P/E', 'P/S', 'Premium / Discount']
    for ci, h in enumerate(comp_hdrs, start=1):
        _set(ws3, 3, ci, h, bold=True, size=9, color=WHITE, fill=NAVY, align='center')

    for ri, c in enumerate(comp or []):
        row = 4 + ri
        prem = str(c.get('premium') or c.get('premDisc') or '')
        is_pos = not prem.startswith('-')
        prem_color = GREEN if is_pos else RED
        _set(ws3, row, 1, c.get('name', ''), size=9, color=OFF_WHITE)
        _set(ws3, row, 2, c.get('ticker', ''), size=9, color=GOLD, align='center')
        if template == 'bank_financials':
            _set(ws3, row, 3, c.get('ptbv', c.get('evEbitda', '')), size=9, color=OFF_WHITE, align='right')
            _set(ws3, row, 4, c.get('pe', ''), size=9, color=OFF_WHITE, align='right')
            _set(ws3, row, 5, c.get('rotce', c.get('ps', '')), size=9, color=OFF_WHITE, align='right')
        else:
            _set(ws3, row, 3, c.get('evEbitda', ''), size=9, color=OFF_WHITE, align='right')
            _set(ws3, row, 4, c.get('pe', ''), size=9, color=OFF_WHITE, align='right')
            _set(ws3, row, 5, c.get('ps', ''), size=9, color=OFF_WHITE, align='right')
        _set(ws3, row, 6, prem, size=9, color=prem_color, align='right')

    # Subject row
    sub_row = 4 + len(comp or []) + 1
    _set(ws3, sub_row, 1, f'{ticker} (Subject)', bold=True, size=9, color=GOLD_LIGHT, fill=NAVY_LIGHT)
    _set(ws3, sub_row, 2, ticker, bold=True, size=9, color=GOLD, fill=NAVY_LIGHT, align='center')
    if template == 'bank_financials':
        _set(ws3, sub_row, 3, km.get('ptbv', ''),   bold=True, size=9, color=GOLD, fill=NAVY_LIGHT, align='right')
        _set(ws3, sub_row, 4, km.get('pe', ''),      bold=True, size=9, color=GOLD, fill=NAVY_LIGHT, align='right')
        _set(ws3, sub_row, 5, km.get('rotce', ''),   bold=True, size=9, color=GOLD, fill=NAVY_LIGHT, align='right')
    else:
        _set(ws3, sub_row, 3, km.get('evEbitda', ''), bold=True, size=9, color=GOLD, fill=NAVY_LIGHT, align='right')
        _set(ws3, sub_row, 4, km.get('pe', ''),       bold=True, size=9, color=GOLD, fill=NAVY_LIGHT, align='right')
        _set(ws3, sub_row, 5, '—',                    bold=True, size=9, color=GOLD, fill=NAVY_LIGHT, align='right')
    _set(ws3, sub_row, 6, '—', bold=True, size=9, color=GOLD, fill=NAVY_LIGHT, align='right')

    ws3.sheet_view.showGridLines = False

    wb.save(str(out_path))


# ── PowerPoint helper ──────────────────────────────────────────────────────────
def _rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _slide_bg(slide, hex_color):
    """Fill slide background with a solid color."""
    bg    = slide.background
    fill  = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex_color)

def _add_textbox(slide, left, top, width, height, text, size=12,
                 bold=False, color=OFF_WHITE, align=PP_ALIGN.LEFT,
                 italic=False, wrap=True):
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    return txBox

def _add_rect(slide, left, top, width, height, fill_hex, line_hex=None):
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = _rgb(line_hex)
    else:
        shape.line.fill.background()
    return shape


# ── PowerPoint Pitch Deck ─────────────────────────────────────────────────────
def generate_pptx(data: dict, ticker: str, today: str, out_path: Path):
    km       = data.get('keyMetrics') or {}
    dcf      = data.get('dcf') or {}
    comp     = data.get('comps') or []
    bull     = data.get('bull') or []
    bear     = data.get('bear') or []
    cats     = data.get('catalysts') or []
    rating   = data.get('rating', 'HOLD')
    target   = data.get('targetPrice', '')
    current  = data.get('currentPrice', '')
    upside   = data.get('upside', '')
    template = select_template(data)
    rating_color = GREEN if rating == 'BUY' else (RED if rating == 'SELL' else AMBER)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank

    def new_slide():
        return prs.slides.add_slide(blank_layout)

    def section_label(slide, text, top=0.28):
        _add_textbox(slide, 0.4, top, 3.0, 0.25, text,
                     size=7, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    def body_text(slide, text, left, top, width, height, size=10):
        _add_textbox(slide, left, top, width, height, text,
                     size=size, color=OFF_WHITE, wrap=True)

    # ── Slide 1: Cover ─────────────────────────────────────────────────────────
    s1 = new_slide()
    _slide_bg(s1, DARK_BG)
    _add_rect(s1, 0, 0, 13.33, 0.08, GOLD)               # gold top stripe
    _add_rect(s1, 0, 7.42, 13.33, 0.08, GOLD)             # gold bottom stripe

    _add_textbox(s1, 0.5, 1.4, 12.0, 1.2, ticker,
                 size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    _add_textbox(s1, 0.5, 2.7, 12.0, 0.6, data.get('company', ''),
                 size=22, bold=False, color=OFF_WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(s1, 0.5, 3.5, 12.0, 0.4, f'QUICK PITCH  ·  {today}',
                 size=10, color=GOLD, align=PP_ALIGN.CENTER)

    # Rating badge
    _add_rect(s1, 5.5, 4.3, 2.33, 0.65, NAVY_LIGHT, GOLD)
    _add_textbox(s1, 5.5, 4.35, 2.33, 0.55,
                 f'{rating}  ·  Target ${target}  ·  {upside}',
                 size=11, bold=True, color=rating_color, align=PP_ALIGN.CENTER)

    _add_textbox(s1, 0.5, 5.3, 12.0, 0.8,
                 f'"{data.get("thesis","")}"',
                 size=10, italic=True, color=OFF_WHITE, align=PP_ALIGN.CENTER, wrap=True)
    _add_textbox(s1, 0.5, 6.9, 12.0, 0.3, 'DJAI Finance Intelligence',
                 size=8, color=GOLD, align=PP_ALIGN.CENTER)

    # ── Slide 2: Business Overview ─────────────────────────────────────────────
    s2 = new_slide()
    _slide_bg(s2, DARK_BG)
    _add_rect(s2, 0, 0, 13.33, 0.08, GOLD)
    _add_textbox(s2, 0.4, 0.2, 12.0, 0.5, 'Business Overview',
                 size=20, bold=True, color=GOLD)
    section_label(s2, f'{ticker}  ·  {data.get("company","")}', top=0.75)
    body_text(s2, data.get('businessOverview', ''), 0.4, 1.1, 12.5, 1.4, size=12)
    section_label(s2, 'KEY SEGMENTS / MOAT', top=2.7)
    body_text(s2, data.get('financialAnalysis', ''), 0.4, 3.0, 12.5, 1.4, size=11)
    section_label(s2, 'RISKS', top=4.55)
    body_text(s2, data.get('risks', ''), 0.4, 4.85, 12.5, 0.9, size=11)

    # ── Slide 3: Financial Snapshot ─────────────────────────────────────────────
    s3 = new_slide()
    _slide_bg(s3, DARK_BG)
    _add_rect(s3, 0, 0, 13.33, 0.08, GOLD)
    _add_textbox(s3, 0.4, 0.2, 12.0, 0.5, 'Financial Snapshot',
                 size=20, bold=True, color=GOLD)

    if template == 'bank_financials':
        metrics_grid = [
            ('Revenue',        km.get('revenue','')),
            ('Rev Growth',     km.get('revenueGrowth','')),
            ('Net Margin',     km.get('netMargin','')),
            ('ROTCE',          km.get('rotce','')),
            ('CET1 Ratio',     km.get('cet1','')),
            ('Efficiency Ratio', km.get('efficiencyRatio','')),
            ('P/E',            km.get('pe','')),
            ('Fwd P/E',        km.get('forwardPe','')),
            ('P/TBV',          km.get('ptbv','')),
            ('Dividend Yield', km.get('dividendYield','')),
            ('Beta',           km.get('beta','')),
            ('Analyst Target', f'${target}' if target else ''),
        ]
    else:
        metrics_grid = [
            ('Revenue',     km.get('revenue','')),
            ('Rev Growth',  km.get('revenueGrowth','')),
            ('Gross Margin',km.get('grossMargin','')),
            ('Op Margin',   km.get('operatingMargin','')),
            ('Net Margin',  km.get('netMargin','')),
            ('FCF',         km.get('fcf','')),
            ('P/E',         km.get('pe','')),
            ('Fwd P/E',     km.get('forwardPe','')),
            ('EV/EBITDA',   km.get('evEbitda','')),
            ('ROE',         km.get('roe','')),
            ('Beta',        km.get('beta','')),
            ('Analyst Target', f'${target}' if target else ''),
        ]
    cols = 4
    cell_w, cell_h = 3.1, 0.75
    for idx, (lbl, val) in enumerate(metrics_grid):
        col = idx % cols
        row = idx // cols
        x = 0.3 + col * (cell_w + 0.12)
        y = 0.95 + row * (cell_h + 0.1)
        _add_rect(s3, x, y, cell_w, cell_h, SURFACE, NAVY_LIGHT)
        _add_textbox(s3, x + 0.1, y + 0.05, cell_w - 0.2, 0.25,
                     lbl.upper(), size=7, bold=True, color=GOLD)
        _add_textbox(s3, x + 0.1, y + 0.32, cell_w - 0.2, 0.35,
                     str(val) if val else '—', size=15, bold=True, color=WHITE)

    # ── Slide 4: Valuation ────────────────────────────────────────────────────
    s4 = new_slide()
    _slide_bg(s4, DARK_BG)
    _add_rect(s4, 0, 0, 13.33, 0.08, GOLD)
    _add_textbox(s4, 0.4, 0.2, 12.0, 0.5, 'Valuation',
                 size=20, bold=True, color=GOLD)

    if template == 'bank_financials':
        val_items = [
            ('P/TBV (Current)',     km.get('ptbv', '—')),
            ('ROTCE',               km.get('rotce', '—')),
            ('CET1 Ratio',          km.get('cet1', '—')),
            ('Efficiency Ratio',    km.get('efficiencyRatio', '—')),
            ('Dividend Yield',      km.get('dividendYield', '—')),
            ('Current Price',       f'${current}'),
            ('Analyst Target',      f'${target}'),
        ]
    else:
        val_items = [
            ('DCF Implied Price',   f'${dcf.get("impliedPrice","—")}'),
            ('WACC',                f'{dcf.get("wacc","")}'),
            ('Terminal Growth',     f'{dcf.get("terminalGrowth","")}'),
            ('Current Price',       f'${current}'),
            ('Upside / (Downside)', upside),
            ('Analyst Target',      f'${target}'),
            ('Rating',              rating),
        ]
    for i, (lbl, val) in enumerate(val_items):
        y = 0.9 + i * 0.7
        _add_rect(s4, 0.4, y, 5.5, 0.55, SURFACE, NAVY_LIGHT)
        _add_textbox(s4, 0.55, y + 0.05, 3.5, 0.45, lbl, size=9, color=GOLD)
        _add_textbox(s4, 4.0, y + 0.05, 1.8, 0.45, str(val), size=11,
                     bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    if template == 'bank_financials':
        section_label(s4, 'VALUATION FRAMEWORK', top=0.9)
        bank_val_note = (
            'Bank Valuation: P/TBV and DDM are the primary frameworks. '
            'Traditional EV/EBITDA or FCF DCF is not applicable. '
            f'Current P/TBV: {km.get("ptbv","—")}. '
            f'ROTCE: {km.get("rotce","—")} drives sustainable book value growth. '
            f'Efficiency ratio: {km.get("efficiencyRatio","—")} (lower = better). '
        )
        body_text(s4, bank_val_note + '\n\n' + data.get('valSummary', ''), 6.5, 0.9, 6.4, 4.5, size=10)
    else:
        section_label(s4, 'VALUATION SUMMARY', top=0.9)
        body_text(s4, data.get('valSummary', ''), 6.5, 0.9, 6.4, 4.5, size=10)

    # ── Slide 5: Bull Case ────────────────────────────────────────────────────
    s5 = new_slide()
    _slide_bg(s5, DARK_BG)
    _add_rect(s5, 0, 0, 13.33, 0.08, GREEN)
    _add_textbox(s5, 0.4, 0.2, 12.0, 0.5, 'Bull Case',
                 size=20, bold=True, color=GREEN)
    _add_textbox(s5, 0.4, 0.8, 12.0, 0.35,
                 f'{ticker}  ·  Why this is a BUY', size=10, color=OFF_WHITE)
    for i, point in enumerate((bull or [])[:3]):
        y = 1.5 + i * 1.5
        _add_rect(s5, 0.4, y, 12.5, 1.2, SURFACE, GREEN)
        _add_textbox(s5, 0.55, y + 0.07, 0.5, 1.0,
                     f'{i+1}.', size=24, bold=True, color=GREEN)
        _add_textbox(s5, 1.1, y + 0.15, 11.4, 0.9, point, size=12, color=WHITE, wrap=True)

    # ── Slide 6: Bear Case ────────────────────────────────────────────────────
    s6 = new_slide()
    _slide_bg(s6, DARK_BG)
    _add_rect(s6, 0, 0, 13.33, 0.08, RED)
    _add_textbox(s6, 0.4, 0.2, 12.0, 0.5, 'Bear Case',
                 size=20, bold=True, color=RED)
    _add_textbox(s6, 0.4, 0.8, 12.0, 0.35,
                 f'{ticker}  ·  Key risks and downside scenarios', size=10, color=OFF_WHITE)
    for i, point in enumerate((bear or [])[:3]):
        y = 1.5 + i * 1.5
        _add_rect(s6, 0.4, y, 12.5, 1.2, SURFACE, RED)
        _add_textbox(s6, 0.55, y + 0.07, 0.5, 1.0,
                     f'{i+1}.', size=24, bold=True, color=RED)
        _add_textbox(s6, 1.1, y + 0.15, 11.4, 0.9, point, size=12, color=WHITE, wrap=True)

    # ── Slide 7: Catalysts & Risks ────────────────────────────────────────────
    s7 = new_slide()
    _slide_bg(s7, DARK_BG)
    _add_rect(s7, 0, 0, 13.33, 0.08, GOLD)
    _add_textbox(s7, 0.4, 0.2, 12.0, 0.5, 'Catalysts & Risks',
                 size=20, bold=True, color=GOLD)

    _add_textbox(s7, 0.4, 0.85, 5.8, 0.3, 'KEY CATALYSTS',
                 size=9, bold=True, color=GREEN)
    for i, cat in enumerate((cats or [])[:3]):
        y = 1.25 + i * 1.1
        _add_rect(s7, 0.4, y, 5.9, 0.85, SURFACE, GREEN)
        _add_textbox(s7, 0.55, y + 0.1, 5.6, 0.7, f'▸  {cat}',
                     size=10, color=WHITE, wrap=True)

    _add_textbox(s7, 7.0, 0.85, 5.8, 0.3, 'KEY RISKS',
                 size=9, bold=True, color=RED)
    risk_text = data.get('risks', '')
    risk_lines = [l.strip() for l in risk_text.split('.') if l.strip()][:3]
    for i, risk in enumerate(risk_lines):
        y = 1.25 + i * 1.1
        _add_rect(s7, 7.0, y, 5.9, 0.85, SURFACE, RED)
        _add_textbox(s7, 7.15, y + 0.1, 5.6, 0.7, f'▾  {risk}.',
                     size=10, color=WHITE, wrap=True)

    # ── Slide 8: Recommendation ───────────────────────────────────────────────
    s8 = new_slide()
    _slide_bg(s8, NAVY)
    _add_rect(s8, 0, 0, 13.33, 0.08, GOLD)
    _add_rect(s8, 0, 7.42, 13.33, 0.08, GOLD)

    _add_textbox(s8, 0.4, 0.3, 12.0, 0.55, 'Recommendation',
                 size=20, bold=True, color=GOLD)
    _add_textbox(s8, 0.4, 1.0, 4.0, 1.5, rating,
                 size=64, bold=True, color=rating_color, align=PP_ALIGN.CENTER)
    _add_textbox(s8, 0.4, 2.6, 4.0, 0.4, f'Target: ${target}  |  {upside} upside',
                 size=12, color=GOLD, align=PP_ALIGN.CENTER)

    _add_rect(s8, 5.0, 0.9, 7.9, 2.2, DARK_BG, GOLD)
    _add_textbox(s8, 5.1, 0.95, 7.7, 0.25, 'INVESTMENT THESIS',
                 size=7, bold=True, color=GOLD)
    _add_textbox(s8, 5.1, 1.25, 7.7, 1.75,
                 f'"{data.get("thesis","")}"',
                 size=11, italic=True, color=OFF_WHITE, wrap=True)

    _add_rect(s8, 0.4, 3.2, 12.5, 1.5, DARK_BG, NAVY_LIGHT)
    _add_textbox(s8, 0.55, 3.25, 12.2, 0.25, 'CONCLUSION',
                 size=7, bold=True, color=GOLD)
    _add_textbox(s8, 0.55, 3.55, 12.2, 1.1,
                 data.get('conclusion', ''), size=10, color=WHITE, wrap=True)

    _add_textbox(s8, 0.4, 5.0, 12.5, 0.3,
                 f'Generated by DJAI Finance Intelligence  ·  {today}  ·  For informational purposes only.',
                 size=7, color=GOLD, align=PP_ALIGN.CENTER)

    prs.save(str(out_path))


# ── Entry point ───────────────────────────────────────────────────────────────
if __name__ == '__main__':
    PITCHES_DIR.mkdir(parents=True, exist_ok=True)
    server = http.server.HTTPServer(('localhost', PORT), ArtifactHandler)
    print(f'[DJAI] Artifact server running on http://localhost:{PORT}')
    print(f'[DJAI] Saving pitches to: {PITCHES_DIR}')
    print(f'[DJAI] openpyxl: {"✓" if HAS_OPENPYXL else "✗ (Excel disabled)"}')
    print(f'[DJAI] pptx:     {"✓" if HAS_PPTX else "✗ (PowerPoint disabled)"}')
    print('[DJAI] Press Ctrl+C to stop.')
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print('\n[DJAI] Server stopped.')
        server.server_close()
